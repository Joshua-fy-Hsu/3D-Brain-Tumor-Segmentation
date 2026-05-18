"""Build the A1 portrait capstone poster for AURAS (PowerPoint, all English).

A1 portrait = 594 x 841 mm = 23.39 x 33.11 in. Single slide.
Figures: docs/report_figures/   Results: results/final/
Output: docs/AURAS_poster_A1.pptx
"""
from pathlib import Path
from PIL import Image, ImageFile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn

ImageFile.LOAD_TRUNCATED_IMAGES = True

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "docs" / "report_figures"
OUT = ROOT / "docs" / "AURAS_poster_A1.pptx"

# ---- palette (AURAS deck conventions) ----
NAVY   = RGBColor(0x16, 0x2A, 0x4A)   # header band / titles
ACCENT = RGBColor(0x3B, 0x6F, 0xB5)   # card title bars
LIGHT  = RGBColor(0xEA, 0xF0, 0xF8)   # card title bar fill (light)
BG     = RGBColor(0xF4, 0xF7, 0xFB)   # slide background
INK    = RGBColor(0x20, 0x2A, 0x38)   # body text
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
NCR_G  = RGBColor(0x22, 0xB3, 0x4D)
ED_B   = RGBColor(0x26, 0x66, 0xD9)
ET_R   = RGBColor(0xDC, 0x29, 0x29)

EMU_IN = 914400
PW, PH = 23.39, 33.11                  # A1 portrait, inches

prs = Presentation()
prs.slide_width  = Emu(int(PW * EMU_IN))
prs.slide_height = Emu(int(PH * EMU_IN))
slide = prs.slides.add_slide(prs.slide_layouts[6])

# slide background
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = BG


def _no_autofit(tf):
    # disable PowerPoint shrink/grow so our sizes are honoured
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    for tag in ('a:normAutofit', 'a:spAutoFit'):
        e = bodyPr.find(qn(tag))
        if e is not None:
            bodyPr.remove(e)


def textbox(x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of dicts {text,size,bold,color,space_after,bullet}"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    _no_autofit(tf)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        p.space_after = Pt(ln.get("space_after", 4))
        p.space_before = Pt(ln.get("space_before", 0))
        if ln.get("bullet"):
            pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
            bu = pPr.makeelement(qn('a:buChar'), {'char': '▪'})
            pPr.append(bu)
            p.level = 0
        r = p.add_run()
        r.text = ln["text"]
        f = r.font
        f.size = Pt(ln.get("size", 18))
        f.bold = ln.get("bold", False)
        f.italic = ln.get("italic", False)
        f.name = "Calibri"
        f.color.rgb = ln.get("color", INK)
    return tb


def rect(x, y, w, h, fill, line=None, line_w=0.75, rounded=False, shadow=False):
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        try:
            shp.adjustments[0] = 0.04
        except Exception:
            pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if shadow:
        el = shp._element.spPr
        sp = el.makeelement(qn('a:effectLst'), {})
        sh = el.makeelement(qn('a:outerShdw'),
                            {'blurRad': '50000', 'dist': '20000',
                             'dir': '5400000', 'rotWithShape': '0'})
        clr = el.makeelement(qn('a:srgbClr'), {'val': '1F2A3A'})
        al = el.makeelement(qn('a:alpha'), {'val': '22000'})
        clr.append(al)
        sh.append(clr)
        sp.append(el.makeelement(qn('a:effectLst'), {}) if False else sh.getparent() or sh)
    return shp


def card(x, y, w, h, title):
    """Card with light title bar; returns inner content rect (x,y,w,h)."""
    rect(x, y, w, h, WHITE, line=RGBColor(0xD3, 0xDD, 0xEA), line_w=1.0, rounded=True)
    bar_h = 0.92
    rect(x, y, w, bar_h, ACCENT, rounded=True)
    rect(x, y + bar_h - 0.18, w, 0.18, ACCENT)  # square off bottom of bar
    textbox(x + 0.28, y, w - 0.5, bar_h,
            [{"text": title, "size": 27, "bold": True, "color": WHITE}],
            anchor=MSO_ANCHOR.MIDDLE)
    pad = 0.28
    return (x + pad, y + bar_h + 0.16, w - 2 * pad, h - bar_h - 0.16 - pad)


def add_fig(path, cx, cy, cw, ch, caption=None):
    """Fit image into (cx,cy,cw,ch) box preserving aspect; centered. Optional caption below."""
    im = Image.open(path)
    iw, ih = im.size
    cap_h = 0.42 if caption else 0.0
    avail_h = ch - cap_h
    ar = iw / ih
    w = cw
    h = w / ar
    if h > avail_h:
        h = avail_h
        w = h * ar
    px = cx + (cw - w) / 2
    py = cy + (avail_h - h) / 2
    slide.shapes.add_picture(str(path), Inches(px), Inches(py),
                             Inches(w), Inches(h))
    if caption:
        textbox(cx, cy + avail_h + 0.02, cw, cap_h,
                [{"text": caption, "size": 14, "italic": True,
                  "color": RGBColor(0x5A, 0x66, 0x78), "align": PP_ALIGN.CENTER}],
                align=PP_ALIGN.CENTER)


# ============================================================ HEADER
M = 0.6                                  # outer margin
HDR_H = 3.55
rect(0, 0, PW, HDR_H, NAVY)
rect(0, HDR_H, PW, 0.10, RGBColor(0x3B, 0x6F, 0xB5))

textbox(M, 0.35, PW - 2 * M, 1.45, [
    {"text": "AURAS — All-modality, Uncertainty-aware, Residual, "
             "Aggregation, Spectral Network", "size": 50, "bold": True,
     "color": WHITE, "align": PP_ALIGN.CENTER}],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
textbox(M, 1.78, PW - 2 * M, 0.8, [
    {"text": "A configurable 3D Trans-ResU-Net for multi-modal brain tumor "
             "segmentation, with built-in uncertainty and a leave-one-out "
             "ablation study", "size": 25, "italic": True,
     "color": RGBColor(0xC9, 0xD8, 0xEC), "align": PP_ALIGN.CENTER}],
    align=PP_ALIGN.CENTER)
textbox(M, 2.66, PW - 2 * M, 0.7, [
    {"text": "Joshua Hsu   •   Brain Tumor Segmentation Capstone   •   "
             "BraTS 2021 (1251 volumes · 1000 train / 251 validation)",
     "size": 20, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}],
    align=PP_ALIGN.CENTER)

# ============================================================ GRID
top = HDR_H + 0.45
foot_h = 1.45
bot = PH - foot_h - 0.35
gut = 0.45
col_w = (PW - 2 * M - 2 * gut) / 3
COLX = [M, M + col_w + gut, M + 2 * (col_w + gut)]


def layout_column(cx, cards):
    """cards: list of (title, weight, render_fn). Heights fill top..bot
    with `gut` between cards, proportional to weight."""
    n = len(cards)
    avail = (bot - top) - gut * (n - 1)
    wsum = sum(w for _, w, _ in cards)
    y = top
    for title, w, fn in cards:
        ch = avail * w / wsum
        ix, iy, iw, ih = card(cx, y, col_w, ch, title)
        fn(ix, iy, iw, ih)
        y += ch + gut


# ---------- renderers ----------
def r_problem(x, y, w, h):
    textbox(x, y, w, h, [
        {"text": "Gliomas are graded from three nested tumor regions, each "
                 "driving a different clinical decision.", "size": 19,
         "space_after": 10},
        {"text": "Whole Tumor (WT) — full extent for surgical planning",
         "size": 19, "bullet": True, "color": ED_B, "bold": True,
         "space_after": 6},
        {"text": "Tumor Core (TC) — resectable mass",
         "size": 19, "bullet": True, "color": NCR_G, "bold": True,
         "space_after": 6},
        {"text": "Enhancing Tumor (ET) — active, aggressive disease",
         "size": 19, "bullet": True, "color": ET_R, "bold": True,
         "space_after": 10},
        {"text": "Manual delineation across 4 MRI modalities is slow, "
                 "subjective and hard to reproduce — and a model should also "
                 "say when it is unsure.", "size": 19},
    ])


def r_data(x, y, w, h):
    add_fig(FIG / "class_imbalance.png", x, y, w, h - 1.95,
            caption="Severe class imbalance — ET is a tiny fraction of "
                    "each volume")
    textbox(x, y + h - 1.8, w, 1.8, [
        {"text": "4 co-registered MRI modalities + foreground mask → "
                 "5-channel 128³ input.", "size": 17, "bullet": True,
         "space_after": 6},
        {"text": "Loss and patch sampling are region-weighted toward the "
                 "rare tumor voxels.", "size": 17, "bullet": True},
    ])


def r_hero(x, y, w, h):
    add_fig(FIG / "conclusion_hero.png", x, y, w, h,
            caption="AURAS prediction vs ground truth — a median-Dice "
                    "validation case")


def r_arch(x, y, w, h):
    add_fig(FIG / "auras_model.png", x, y, w, h - 2.35)
    textbox(x, y + h - 2.2, w, 2.2, [
        {"text": "One configurable 3D Trans-ResU-Net backbone.", "size": 18,
         "bold": True, "space_after": 6},
        {"text": "Per-modality stems + cross-modal attention fuse the 4 MRIs",
         "size": 17, "bullet": True, "space_after": 5},
        {"text": "Frequency block + spectral-Swin stage add global context",
         "size": 17, "bullet": True, "space_after": 5},
        {"text": "Predictive-variance head emits a per-voxel uncertainty map",
         "size": 17, "bullet": True, "space_after": 5},
        {"text": "Boundary head sharpens tumor edges", "size": 17,
         "bullet": True},
    ])


def r_design(x, y, w, h):
    add_fig(FIG / "roadmap.png", x, y, w, h - 1.7)
    textbox(x, y + h - 1.6, w, 1.6, [
        {"text": "Components are stacked one at a time; a true "
                 "leave-one-out test then removes exactly one piece from "
                 "the full model — so every gain is attributable.",
         "size": 17},
    ])


def r_loss(x, y, w, h):
    add_fig(FIG / "Loss_Function.jpg", x, y, w, h - 2.6)
    textbox(x, y + h - 2.5, w, 2.5, [
        {"text": "Region-wise Dice + Focal loss on WT / TC / ET with deep "
                 "supervision.", "size": 17, "bullet": True,
         "space_after": 6},
        {"text": "Uncertainty- and boundary-aware terms ramped in over "
                 "training.", "size": 17, "bullet": True, "space_after": 6},
        {"text": "EMA weights, bf16, top-5 snapshot ensemble + flip-TTA at "
                 "inference.", "size": 17, "bullet": True},
    ])


def r_results(x, y, w, h):
    rows = [
        ("Region", "Base CNN", "AURAS", "Δ"),
        ("Dice WT", "0.923", "0.929", "+0.6"),
        ("Dice TC", "0.788", "0.781", "−0.7"),
        ("Dice ET", "0.765", "0.789", "+2.4"),
        ("HD95 (mm)", "6.75", "5.95", "−0.80"),
        ("NSD", "0.787", "0.803", "+1.6"),
    ]
    t_h = 3.3
    tb = slide.shapes.add_table(len(rows), 4, Inches(x), Inches(y),
                                Inches(w), Inches(t_h)).table
    tb.first_row = False
    for j, fr in enumerate([0.30, 0.25, 0.23, 0.22]):
        tb.columns[j].width = Inches(w * fr)
    for i, row in enumerate(rows):
        tb.rows[i].height = Inches(t_h / len(rows))
        for j, val in enumerate(row):
            cl = tb.cell(i, j)
            cl.text = val
            pr = cl.text_frame.paragraphs[0]
            pr.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            rn = pr.runs[0]
            rn.font.size = Pt(17)
            rn.font.name = "Calibri"
            rn.font.bold = (i == 0) or (j == 0)
            rn.font.color.rgb = WHITE if i == 0 else INK
            tcPr = cl._tc.get_or_add_tcPr()
            fillel = tcPr.makeelement(qn('a:solidFill'), {})
            clr = tcPr.makeelement(qn('a:srgbClr'),
                                   {'val': '3B6FB5' if i == 0 else
                                    ('EDF2F8' if i % 2 else 'FFFFFF')})
            fillel.append(clr)
            tcPr.append(fillel)
            cl.margin_left = cl.margin_right = Pt(5)
            cl.margin_top = cl.margin_bottom = Pt(2)
    textbox(x, y + t_h + 0.2, w, h - t_h - 0.2, [
        {"text": "AURAS improves WT Dice and HD95 over the Base CNN with "
                 "statistical significance (paired Wilcoxon, Bonferroni "
                 "p < 0.01) and lifts ET Dice by ~2.4 points, at equal "
                 "specificity.", "size": 17}])


def r_ablation(x, y, w, h):
    add_fig(FIG / "ablation_dice.png", x, y, w, h,
            caption="Per-region Dice as each component is added to the "
                    "backbone")


def r_cost(x, y, w, h):
    add_fig(FIG / "complexity_tradeoff.png", x, y, w, h - 1.1,
            caption="37 M params · ~85 ms / volume · ~2.4 GB VRAM at 128³")


def r_deploy(x, y, w, h):
    add_fig(FIG / "workstation.png", x, y, w, h,
            caption="Web workstation — segmentation with per-region "
                    "confidence")


layout_column(COLX[0], [
    ("1.  Clinical Problem", 1.05, r_problem),
    ("2.  Data & Challenge", 1.05, r_data),
    ("Qualitative Result",   0.95, r_hero),
])
layout_column(COLX[1], [
    ("3.  AURAS Architecture", 1.55, r_arch),
    ("4.  Ablation Design",    0.80, r_design),
    ("5.  Training & Loss",    0.95, r_loss),
])
layout_column(COLX[2], [
    ("6.  Headline Results",  0.95, r_results),
    ("7.  Component Ablation", 0.85, r_ablation),
    ("8.  Accuracy vs Cost",   0.78, r_cost),
    ("9.  Deployment",         0.92, r_deploy),
])

# ============================================================ FOOTER
fy = PH - foot_h
rect(0, fy, PW, foot_h, NAVY)
textbox(M, fy + 0.12, PW - 2 * M, foot_h - 0.2, [
    {"text": "Conclusion", "size": 21, "bold": True, "color": WHITE,
     "align": PP_ALIGN.CENTER, "space_after": 3},
    {"text": "A single configurable backbone, validated by true "
             "leave-one-out ablation, delivers calibrated, uncertainty-aware "
             "tumor segmentation that significantly improves whole-tumor "
             "accuracy and boundary quality — deployed in an interactive "
             "clinical workstation.", "size": 18,
     "color": RGBColor(0xDD, 0xE6, 0xF3), "align": PP_ALIGN.CENTER}],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

prs.save(str(OUT))
print("saved", OUT)
